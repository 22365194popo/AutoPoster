#!/usr/bin/env python
# -*- coding: utf-8 -*-
# @Author  : Jay
# @Time    : 2020/5/12

import sys
import os
import cv2
import numpy as np
import re
import datetime
import pickle
import json
import pptx
from PIL import ImageFont, ImageDraw, Image
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt

TIME = '13:00 - 15:00'
BUILDING = '格致大樓'
HEIGHT = [230, 330, 370, 480, 570, 610]
FONT_PATH = 'font/simsun.ttc'
FONT_SIZE = 32
DICT_PATH = './sources/result_file/speech.pkl'
IMG_PATH = './sources/result_img/poster.png'
ADD_IMG_PATH = './sources/result_img/'
ADD_FILE_PATH = './sources/result_file/'

def show(d):
    ''' Generate a speech poster.

    Args:
        d (dict): The dict is speecher data from app.py.

    Examples:
        >>> poster.show(dictionary)
        Image will be download in your IMG_PATH.
    '''
    day = eval_date(d)
    sentences = [d['演講題目'], d['演講者'], d['學歷/公司'], BUILDING, day, TIME]
    bk_img = cv2.imread(IMG_PATH)
    font = ImageFont.truetype(FONT_PATH, FONT_SIZE)
    img_pil = Image.fromarray(bk_img)
    draw = ImageDraw.Draw(img_pil)
    img = Image.open(IMG_PATH)

    for i in range(len(sentences)):
        num_not_chinese = 0
        if re.findall('[\x00-\xff]', sentences[i]) is not None:
            num_not_chinese = len(re.findall('[\x00-\xff]', sentences[i]))
            initial_point = img.size[0]/2 - (32*(len(sentences[i])-num_not_chinese/2)/2)
        else:
            initial_point = img.size[0]/2 - (32*len(sentences[i]/2))
        draw.text((initial_point, HEIGHT[i]), sentences[i], font=font, fill=(0, 0, 0))

    bk_img = np.array(img_pil)
    cv2.imshow(sentences[0], bk_img) # show img.
    # press any key to shut down window.
    cv2.waitKey(0)
    cv2.destroyAllWindows() 
    cv2.imwrite(ADD_IMG_PATH + sentences[0] + '.png', bk_img)

def eval_font():
    ''' To check if there are oversize sentence in image. '''
    #
        #if sentence >= 整個圖片寬度的百分之60: # 不美觀 須調整字體
        #    長度 <- 整個圖片寬度的百分之60
        #    字體大小 <- int(長度 / 有幾個字)
    
    for sentence in lines:
        if len(sentence) >= img.size[0] * 0.6:
            FONT_SIZE = int(img.size[0]*0.6 / len(sentence))
            break

def eval_height():
    ''' According to the height of title, we calcualate the difference of height of each title.'''
    title_height = [180, 280, 450, 535]
    title_gap = [(title_height[i+1] - title_height[i]) for i in range(len(title_height))]


def eval_date(d):
    ''' In order to get the weekday. 
    
    Args:
        d (dict): The dict is speecher data from app.py.

    return:
           (str): The date will be posted in poster.

    Examples:
        >>> eval_date(d)
        return date
    '''
    date = d['演講時間']
    mat = re.search(r"(\d{4}-\d{1,2}-\d{1,2})",date)
    rep_mat = mat.group(0).replace('-', '') # 2020-12-20 -> 20201220
    spt_date = mat.group(0).split('-') # 2020-12-20 -> ['2020', '12', '20']
    week = datetime.strptime(rep_mat, '%Y%m%d').weekday() # 0~6 for 星期一 ~ 星期日

    if week == '0':
        day = '一'
    elif week == '1':
        day = '二'
    elif week == '2':
        day = '三'
    elif week == '3':
        day = '四'
    elif week == '4':
        day = '五'
    elif week == '5':
        day = '六'
    else:
        day = '日'

    return (spt_date[0]  + '年' + spt_date[1]  + '月' + spt_date[2]  + '日' + '星期' + day) 

def pptx(d):
    ''' Using pptx modules to change .jpg file to .pptx file. '''
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # type six is for empty slide.
    shapes = slide.shapes
    picture = slide.shapes.add_picture(ADD_IMG_PATH + 'poster' + '.png', 25.4, 14.288) # add image.

    textbox_title = slide.shapes.add_textbox(0.5, 5, 5, 5)  # left，top爲相對位置，width，height爲文本框大小
    textbox_speecher = slide.shapes.add_textbox(0.15, 15, 10, 10)
    textbox_building = slide.shapes.add_textbox(0.25, 25, 15, 15)
    textbox_date = slide.shapes.add_textbox(0.35, 35, 20, 20)

    textbox_date.text = d['演講題目']
    textbox_speecher.text = d['演講者']
    textbox_building.text = BUILDING
    textbox_date.text = eval_date(d)

    

    prs.save(ADD_FILE_PATH + d['演講題目'] + '.pptx')

def write(d):
    ''' Write data.pkl in folder. '''
    f = open(DICT_PATH, 'wb')
    pickle.dump(d, f)
    f.close()

def read():
    ''' Read data.pkl from folder. '''
    f = open(DICT_PATH, 'rb')
    output = pickle.load(f)
    print(output)